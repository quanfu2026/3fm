"""
?????? - ???????
???Markdown / PDF / DOCX / XLSX / PPTX / ??
"""
import os
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

KB_PATH   = os.getenv("KNOWLEDGE_BASE_PATH", "./knowledge_base")
DOCS_PATH = os.path.join(KB_PATH, "docs")
CHUNK_SIZE   = 300
CHUNK_OVERLAP = 50


class KnowledgeBuilder:

    def load_all_documents(self) -> list:
        docs  = []
        path  = Path(DOCS_PATH)
        path.mkdir(parents=True, exist_ok=True)

        for file in sorted(path.rglob("*")):
            if file.is_dir():
                continue
            try:
                chunks = self._process_file(file)
                docs.extend(chunks)
                logger.info(f"  {file.name}?{len(chunks)} chunks")
            except Exception as e:
                logger.warning(f"  {file.name}?{e}")

        logger.info(f"??????{len(docs)} chunks")
        print("[??] ?? docs?", len(docs))
        return docs

    def _process_file(self, file: Path) -> list:
        ext = file.suffix.lower()
        if ext == ".md":
            return self._from_markdown(file)
        elif ext == ".pdf":
            return self._from_pdf(file)
        elif ext == ".docx":
            return self._from_docx(file)
        elif ext in (".xlsx", ".xls"):
            return self._from_xlsx(file)
        elif ext == ".pptx":
            return self._from_pptx(file)
        elif ext in (".jpg", ".jpeg", ".png", ".webp"):
            return self._from_image(file)
        elif ext == ".txt":
            return self._from_text(file)
        else:
            return []

    def _from_markdown(self, file: Path) -> list:
        text = file.read_text(encoding="utf-8")
        print("[??] markdown ??", len(text), "?")
        chunks = self._chunk(text, str(file), "markdown")
        print("[??] ??", len(chunks), "chunks")
        return chunks

    def _from_text(self, file: Path) -> list:
        text = file.read_text(encoding="utf-8")
        return self._chunk(text, str(file), "text")

    def _from_pdf(self, file: Path) -> list:
        import fitz
        doc  = fitz.open(str(file))
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return self._chunk(text, str(file), "pdf")

    def _from_docx(self, file: Path) -> list:
        from docx import Document
        doc  = Document(str(file))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return self._chunk(text, str(file), "docx")

    def _from_xlsx(self, file: Path) -> list:
        import openpyxl
        wb    = openpyxl.load_workbook(str(file), read_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"?????{ws.title}?")
            for row in ws.iter_rows(values_only=True):
                row_str = " | ".join(str(c) for c in row if c is not None)
                if row_str.strip():
                    lines.append(row_str)
        wb.close()
        return self._chunk("\n".join(lines), str(file), "xlsx")

    def _from_pptx(self, file: Path) -> list:
        from pptx import Presentation
        prs   = Presentation(str(file))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"?? {i} ?????")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
        return self._chunk("\n".join(lines), str(file), "pptx")

    def _from_image(self, file: Path) -> list:
        try:
            import pytesseract
            from PIL import Image
            img  = Image.open(str(file))
            text = pytesseract.image_to_string(img, lang="chi_tra+eng")
            if text.strip():
                return self._chunk(text, str(file), "image_ocr")
        except Exception:
            pass
        text = f"?????{file.name}????{file}?"
        return self._chunk(text, str(file), "image")

    def _chunk(self, text: str, source: str, doc_type: str) -> list:
        text   = text.strip()
        if not text:
            return []
        chunks = []
        start  = 0
        idx    = 0
        fname  = Path(source).stem
        while start < len(text):
            end   = min(start + CHUNK_SIZE, len(text))
            chunk = text[start:end].strip()
            if chunk:
                chunks.append({
                    "id":      f"{fname}_{idx}",
                    "content": chunk,
                    "source":  source,
                    "type":    doc_type,
                })
                idx += 1
            if end >= len(text):
                break
            start = end - CHUNK_OVERLAP
        return chunks
