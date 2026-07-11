"""
Knowledge Builder
支援 Markdown、TXT、PDF、DOCX、XLSX、PPTX 與圖片。

切塊策略：
1. FAQ 文件：一組 Q&A 為一個 Chunk。
2. 一般文件：使用固定長度與重疊區間切塊。
"""

import os
import re
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

KB_PATH = os.getenv("KNOWLEDGE_BASE_PATH", "./knowledge_base")
DOCS_PATH = os.path.join(KB_PATH, "docs")

CHUNK_SIZE = 500
CHUNK_OVERLAP = 80


class KnowledgeBuilder:

    def load_all_documents(self) -> list:
        docs = []
        path = Path(DOCS_PATH)
        path.mkdir(parents=True, exist_ok=True)

        for file in sorted(path.rglob("*")):
            if file.is_dir():
                continue

            try:
                chunks = self._process_file(file)
                docs.extend(chunks)
                logger.info(
                    "%s：%s chunks",
                    file.name,
                    len(chunks),
                )
            except Exception as exc:
                logger.warning(
                    "%s 處理失敗：%s",
                    file.name,
                    exc,
                )

        logger.info("知識庫共建立 %s chunks", len(docs))
        print("[KB] 文件切塊總數：", len(docs))
        return docs

    def _process_file(self, file: Path) -> list:
        ext = file.suffix.lower()

        if ext == ".md":
            return self._from_markdown(file)

        if ext == ".txt":
            return self._from_text(file)

        if ext == ".pdf":
            return self._from_pdf(file)

        if ext == ".docx":
            return self._from_docx(file)

        if ext in {".xlsx", ".xls"}:
            return self._from_xlsx(file)

        if ext == ".pptx":
            return self._from_pptx(file)

        if ext in {".jpg", ".jpeg", ".png", ".webp"}:
            return self._from_image(file)

        return []

    def _from_markdown(self, file: Path) -> list:
        text = file.read_text(
            encoding="utf-8",
            errors="ignore",
        )

        print(
            "[KB] Markdown：",
            file.name,
            len(text),
            "字元",
        )

        chunks = self._smart_chunk(
            text=text,
            source=str(file),
            doc_type="markdown",
        )

        print(
            "[KB] 建立",
            len(chunks),
            "chunks",
        )
        return chunks

    def _from_text(self, file: Path) -> list:
        text = file.read_text(
            encoding="utf-8",
            errors="ignore",
        )

        return self._smart_chunk(
            text=text,
            source=str(file),
            doc_type="text",
        )

    def _from_pdf(self, file: Path) -> list:
        import fitz

        document = fitz.open(str(file))

        try:
            text = "\n".join(
                page.get_text()
                for page in document
            )
        finally:
            document.close()

        return self._smart_chunk(
            text=text,
            source=str(file),
            doc_type="pdf",
        )

    def _from_docx(self, file: Path) -> list:
        from docx import Document

        document = Document(str(file))

        text = "\n".join(
            paragraph.text
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        )

        return self._smart_chunk(
            text=text,
            source=str(file),
            doc_type="docx",
        )

    def _from_xlsx(self, file: Path) -> list:
        import openpyxl

        workbook = openpyxl.load_workbook(
            str(file),
            read_only=True,
            data_only=True,
        )

        lines = []

        try:
            for worksheet in workbook.worksheets:
                lines.append(
                    f"工作表：{worksheet.title}"
                )

                for row in worksheet.iter_rows(
                    values_only=True
                ):
                    row_text = " | ".join(
                        str(cell)
                        for cell in row
                        if cell is not None
                    )

                    if row_text.strip():
                        lines.append(row_text)
        finally:
            workbook.close()

        return self._smart_chunk(
            text="\n".join(lines),
            source=str(file),
            doc_type="xlsx",
        )

    def _from_pptx(self, file: Path) -> list:
        from pptx import Presentation

        presentation = Presentation(str(file))
        lines = []

        for slide_index, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            lines.append(f"投影片 {slide_index}")

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()

                    if text:
                        lines.append(text)

        return self._smart_chunk(
            text="\n".join(lines),
            source=str(file),
            doc_type="pptx",
        )

    def _from_image(self, file: Path) -> list:
        try:
            import pytesseract
            from PIL import Image

            image = Image.open(str(file))

            text = pytesseract.image_to_string(
                image,
                lang="chi_tra+eng",
            )

            if text.strip():
                return self._smart_chunk(
                    text=text,
                    source=str(file),
                    doc_type="image_ocr",
                )

        except Exception as exc:
            logger.warning(
                "圖片 OCR 失敗：%s",
                exc,
            )

        text = (
            f"圖片檔案：{file.name}\n"
            f"來源：{file}"
        )

        return self._fixed_chunk(
            text=text,
            source=str(file),
            doc_type="image",
        )

    def _smart_chunk(
        self,
        text: str,
        source: str,
        doc_type: str,
    ) -> list:
        """
        優先嘗試 FAQ 切塊。

        若至少找到兩組 Q&A，就採用：
        一組問題＋答案＝一個 Chunk。

        否則使用一般固定長度切塊。
        """
        normalized = self._normalize_text(text)

        faq_chunks = self._faq_chunk(
            text=normalized,
            source=source,
            doc_type=doc_type,
        )

        if len(faq_chunks) >= 2:
            print(
                f"[KB] FAQ 模式：{Path(source).name}，"
                f"{len(faq_chunks)} 組問答"
            )
            return faq_chunks

        return self._fixed_chunk(
            text=normalized,
            source=source,
            doc_type=doc_type,
        )

    @staticmethod
    def _normalize_text(text: str) -> str:
        text = text.replace("\r\n", "\n")
        text = text.replace("\r", "\n")
        text = text.replace("\u00a0", " ")

        # 保留段落，但移除過多空白行。
        text = re.sub(r"\n{3,}", "\n\n", text)

        return text.strip()

    def _faq_chunk(
        self,
        text: str,
        source: str,
        doc_type: str,
    ) -> list:
        """
        支援下列格式：

        Q1
        Q：多久收到商品？
        A：約 2～5 個工作天。

        Q：多久收到商品？
        A：約 2～5 個工作天。

        Q: How long?
        A: 2-5 business days.
        """

        pattern = re.compile(
            r"""
            (?:
                ^|\n
            )
            \s*
            (?:
                Q\s*\d+\s*
                (?:\n+)?
            )?
            Q\s*[：:]\s*
            (?P<question>.+?)
            \s*
            (?:\n+|\s+)
            A\s*[：:]\s*
            (?P<answer>.+?)
            (?=
                \n\s*
                (?:
                    Q\s*\d+\s*
                    (?:\n+)?
                )?
                Q\s*[：:]
                |
                \Z
            )
            """,
            re.IGNORECASE
            | re.MULTILINE
            | re.DOTALL
            | re.VERBOSE,
        )

        matches = list(pattern.finditer(text))

        if not matches:
            return []

        fname = Path(source).stem
        chunks = []

        for index, match in enumerate(matches):
            question = self._clean_inline(
                match.group("question")
            )
            answer = self._clean_inline(
                match.group("answer")
            )

            if not question or not answer:
                continue

            content = (
                f"問題：{question}\n"
                f"答案：{answer}"
            )

            chunks.append({
                "id": f"{fname}_faq_{index}",
                "content": content,
                "source": source,
                "type": f"{doc_type}_faq",
                "question": question,
                "answer": answer,
            })

        return chunks

    @staticmethod
    def _clean_inline(text: str) -> str:
        text = text.strip()
        text = re.sub(r"\s*\n\s*", " ", text)
        text = re.sub(r"[ \t]{2,}", " ", text)
        return text.strip()

    def _fixed_chunk(
        self,
        text: str,
        source: str,
        doc_type: str,
    ) -> list:
        text = text.strip()

        if not text:
            return []

        chunks = []
        start = 0
        index = 0
        fname = Path(source).stem

        while start < len(text):
            end = min(
                start + CHUNK_SIZE,
                len(text),
            )

            # 優先在句號、換行處結束，避免切斷句子。
            if end < len(text):
                candidate_end = max(
                    text.rfind("\n", start, end),
                    text.rfind("。", start, end),
                    text.rfind("！", start, end),
                    text.rfind("？", start, end),
                )

                if candidate_end > start + 100:
                    end = candidate_end + 1

            chunk = text[start:end].strip()

            if chunk:
                chunks.append({
                    "id": f"{fname}_{index}",
                    "content": chunk,
                    "source": source,
                    "type": doc_type,
                })
                index += 1

            if end >= len(text):
                break

            next_start = end - CHUNK_OVERLAP

            # 防止異常造成無限迴圈。
            if next_start <= start:
                next_start = end

            start = next_start

        return chunks
